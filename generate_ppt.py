import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_path):
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_BG_DARK = RGBColor(15, 23, 42)       # #0F172A
    COLOR_CARD_BG = RGBColor(30, 41, 59)       # #1E293B
    COLOR_CARD_BORDER = RGBColor(51, 65, 85)   # #334155
    COLOR_PRIMARY = RGBColor(14, 165, 233)     # #0EA5E9 (Cyan/Blue)
    COLOR_ACCENT = RGBColor(99, 102, 241)      # #6366F1 (Indigo/Purple)
    COLOR_SUCCESS = RGBColor(16, 185, 129)     # #10B981 (Emerald)
    COLOR_WARNING = RGBColor(245, 158, 11)     # #F59E0B (Amber)
    COLOR_TEXT_LIGHT = RGBColor(248, 250, 252) # #F8FAFC (White/Light)
    COLOR_TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8 (Slate 400)
    COLOR_HIGHLIGHT = RGBColor(56, 189, 248)   # #38BDF8 (Sky Blue)

    def set_slide_background(slide, color=COLOR_BG_DARK):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="AI PATH FINDER PROJECT"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(10)
        p_c.font.bold = True
        p_c.font.color.rgb = COLOR_PRIMARY

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.6))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_LIGHT

        # Subtle divider line
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.35), Inches(11.733), Inches(0.02)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD_BORDER
        shape.line.fill.background()

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1)
        else:
            card.line.fill.background()
        return card

    # ==========================================
    # SLIDE 1: Title Slide (Cover)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)

    # Decorative Card in the center
    add_card(slide1, Inches(1.2), Inches(1.2), Inches(10.933), Inches(5.1), bg_color=COLOR_CARD_BG, border_color=COLOR_PRIMARY)

    # Pill badge
    badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(1.7), Inches(3.2), Inches(0.38))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(30, 58, 138)
    badge.line.color.rgb = COLOR_PRIMARY
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    p_b.text = "NEXT-GEN ADAPTIVE EDTECH"
    p_b.font.size = Pt(10)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_HIGHLIGHT

    # Project Title
    t_box = slide1.shapes.add_textbox(Inches(1.8), Inches(2.3), Inches(9.8), Inches(1.3))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI Path Finder"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_LIGHT

    # Subtitle
    sub_box = slide1.shapes.add_textbox(Inches(1.8), Inches(3.6), Inches(9.8), Inches(0.9))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Personalized Learning Path Recommender with Real-Time Adaptation & AI Mentorship"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = COLOR_PRIMARY

    desc_box = slide1.shapes.add_textbox(Inches(1.8), Inches(4.5), Inches(9.8), Inches(0.8))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "Combines NLP profile extraction, 8-factor recommendation algorithms, adaptive feedback loops, and LLM-powered context mentorship into an end-to-end learning platform."
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = COLOR_TEXT_MUTED

    # Footer / Author in cover
    f_box = slide1.shapes.add_textbox(Inches(1.8), Inches(5.5), Inches(9.8), Inches(0.4))
    tf_f = f_box.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "Full-Stack Project: FastAPI + React TypeScript + Google Gemini AI + SQLite"
    p_f.font.size = Pt(11)
    p_f.font.bold = True
    p_f.font.color.rgb = COLOR_SUCCESS

    # ==========================================
    # SLIDE 2: Problem & Solution
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "The Challenge vs. The Solution", "PROBLEM & SOLUTION STATEMENT")

    # Card 1: Problem
    add_card(slide2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), bg_color=RGBColor(30, 27, 46), border_color=RGBColor(239, 68, 68))
    p_header = slide2.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.0), Inches(0.5))
    tf = p_header.text_frame
    p = tf.paragraphs[0]
    p.text = "❌ The Modern Learning Crisis"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(248, 113, 113)

    p_body = slide2.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(5.0), Inches(4.0))
    tf = p_body.text_frame
    tf.word_wrap = True
    problems = [
        ("Information Overload", "Learners are overwhelmed by thousands of disconnected courses, videos, and tutorials without clear structure."),
        ("One-Size-Fits-All Roadmaps", "Static curriculums ignore individual prior knowledge, unique career goals, and realistic weekly time commitments."),
        ("Static & Non-Adaptive", "Traditional platforms do not adapt when content is too difficult, too easy, or when learner goals evolve."),
        ("No Active Feedback Guidance", "Learners get stuck without contextual mentorship or clear insight into their real skill gaps.")
    ]
    for i, (title, text) in enumerate(problems):
        p1 = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p1.text = f"• {title}:"
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_LIGHT
        p2 = tf.add_paragraph()
        p2.text = f"  {text}\n"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # Card 2: Solution
    add_card(slide2, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2), bg_color=RGBColor(17, 34, 51), border_color=COLOR_SUCCESS)
    s_header = slide2.shapes.add_textbox(Inches(7.2), Inches(1.9), Inches(5.0), Inches(0.5))
    tf = s_header.text_frame
    p = tf.paragraphs[0]
    p.text = "✨ The AI Path Finder Solution"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUCCESS

    s_body = slide2.shapes.add_textbox(Inches(7.2), Inches(2.6), Inches(5.0), Inches(4.0))
    tf = s_body.text_frame
    tf.word_wrap = True
    solutions = [
        ("Natural Language Profiling", "Learners express background & career targets in plain English; AI parses precise skills and timelines."),
        ("Dynamic 8-Factor Roadmaps", "Tailored phased roadmaps that calculate difficulty, prerequisites, and milestone projects."),
        ("Real-Time Adaptive Engine", "Instant curriculum restructuring on learner feedback (Too Hard / Too Easy / Skip / Pace)."),
        ("24/7 Context-Grounded AI Mentor", "An LLM assistant that knows the learner's exact progress, roadmaps, and weak skills.")
    ]
    for i, (title, text) in enumerate(solutions):
        p1 = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p1.text = f"• {title}:"
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_LIGHT
        p2 = tf.add_paragraph()
        p2.text = f"  {text}\n"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 3: Key Features & Capabilities
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "Core Features & System Capabilities", "FEATURE MATRIX")

    features = [
        ("🎯 NLP Profile Extraction", "Extracts target roles, current skill levels, weekly bandwidth, and target completion dates from natural language.", COLOR_PRIMARY),
        ("📊 Skill Gap Radar Analysis", "Compares learner capabilities against industry target role benchmarks with priority rankings.", COLOR_ACCENT),
        ("🗺️ Phased Smart Roadmaps", "Multi-stage milestones, hands-on projects, estimated hours, and prerequisite dependencies.", COLOR_HIGHLIGHT),
        ("🧠 24/7 AI Mentor", "Context-aware chatbot with full access to the learner's live progress, goals, and learning bottlenecks.", COLOR_PRIMARY),
        ("⚡ Dynamic Adaptation", "Real-time adjustments when users rate modules (too easy/hard), adding prerequisites or skipping steps.", COLOR_SUCCESS),
        ("📝 AI Quizzes & Mastery", "Automated skill assessments with immediate feedback and automatic proficiency score updates.", COLOR_WARNING),
    ]

    for idx, (ftitle, fdesc, color) in enumerate(features):
        row = idx // 3
        col = idx % 3
        c_left = Inches(0.8 + col * 4.0)
        c_top = Inches(1.6 + row * 2.6)
        
        card = add_card(slide3, c_left, c_top, Inches(3.7), Inches(2.35), bg_color=COLOR_CARD_BG, border_color=color)
        
        tb = slide3.shapes.add_textbox(c_left + Inches(0.2), c_top + Inches(0.2), Inches(3.3), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p_t = tf.paragraphs[0]
        p_t.text = ftitle
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        
        p_d = tf.add_paragraph()
        p_d.text = f"\n{fdesc}"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 4: System Architecture
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "End-to-End System Architecture", "SYSTEM DESIGN & DATA FLOW")

    arch_layers = [
        ("1. Presentation Layer (Frontend)", "React 19 + TypeScript + Vite\n• Responsive Dashboard, Radar Chart, Roadmap Visualizer\n• Interactive AI Mentor Chat & Assessment Studio\n• State synchronization via RESTful API clients", COLOR_PRIMARY),
        ("2. API & Application Gateway", "FastAPI (Python 3.10+)\n• High-performance asynchronous endpoints\n• Pydantic validation schemas & request pipelines\n• CORS middleware & clean modular routing structure", COLOR_ACCENT),
        ("3. Intelligence & Adaptation Core", "Hybrid AI & Recommendation Engines\n• Google Gemini AI integration (with deterministic offline fallback)\n• 8-Factor scoring matrix & heuristic gap analyzer\n• Real-time adaptation feedback handler", COLOR_SUCCESS),
        ("4. Persistence & Data Layer", "SQLAlchemy ORM + SQLite / AsyncIO\n• Learners, Skill Gaps, Roadmap Items, Feedback & Quiz models\n• Seed data engine with Alex Chen pre-loaded demo profile", COLOR_WARNING)
    ]

    for idx, (ltitle, ldesc, lcolor) in enumerate(arch_layers):
        top_pos = Inches(1.6 + idx * 1.3)
        add_card(slide4, Inches(0.8), top_pos, Inches(11.733), Inches(1.15), bg_color=COLOR_CARD_BG, border_color=lcolor)
        
        tb = slide4.shapes.add_textbox(Inches(1.1), top_pos + Inches(0.12), Inches(11.1), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = ltitle
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = lcolor
        
        p1 = tf.add_paragraph()
        p1.text = ldesc
        p1.font.size = Pt(10.5)
        p1.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 5: Recommendation Algorithm & Scoring Matrix
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "8-Factor Recommendation Scoring Matrix", "INTELLIGENT ALGORITHM")

    # Left: Explanation Card
    add_card(slide5, Inches(0.8), Inches(1.6), Inches(4.5), Inches(5.2), bg_color=COLOR_CARD_BG, border_color=COLOR_PRIMARY)
    tb_left = slide5.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(3.9), Inches(4.8))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "How the Algorithm Works"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_HIGHLIGHT

    p2 = tf_l.add_paragraph()
    p2.text = (
        "\nThe engine calculates a dynamic 'Next Best Action' recommendation for every learner."
        "\n\nEach candidate topic and project is evaluated against 8 weighted multi-dimensional variables."
        "\n\nHigh confidence recommendations prioritize items with prerequisite fulfillment and critical skill gaps."
    )
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = COLOR_TEXT_MUTED

    # Right: Weights Table Cards
    weights = [
        ("Goal Alignment", "25%", "Matches target career role requirements", COLOR_PRIMARY),
        ("Skill Gap Criticality", "20%", "Focuses on largest deficiency areas first", COLOR_ACCENT),
        ("Prerequisites Met", "15%", "Ensures foundational knowledge is mastered", COLOR_SUCCESS),
        ("Difficulty Match", "10%", "Calibrates challenge to learner's current level", COLOR_WARNING),
        ("Learning Style Fit", "10%", "Adapts to visual, practical, or theoretical preference", COLOR_HIGHLIGHT),
        ("Time Availability", "10%", "Aligns estimated item hours with weekly bandwidth", COLOR_PRIMARY),
        ("Progress Momentum", "5%", "Rewards consistency and recent completions", COLOR_SUCCESS),
        ("Feedback Signals", "5%", "Incorporates learner pacing & difficulty ratings", COLOR_ACCENT)
    ]

    for idx, (factor, weight, desc, color) in enumerate(weights):
        r = idx // 2
        c = idx % 2
        card_l = Inches(5.6 + c * 3.5)
        card_t = Inches(1.6 + r * 1.25)
        
        add_card(slide5, card_l, card_t, Inches(3.35), Inches(1.1), bg_color=RGBColor(24, 32, 47), border_color=color)
        
        tb = slide5.shapes.add_textbox(card_l + Inches(0.15), card_t + Inches(0.1), Inches(3.05), Inches(0.9))
        tf_w = tb.text_frame
        tf_w.word_wrap = True
        
        p_w0 = tf_w.paragraphs[0]
        p_w0.text = f"{factor}  •  {weight}"
        p_w0.font.size = Pt(12)
        p_w0.font.bold = True
        p_w0.font.color.rgb = color
        
        p_w1 = tf_w.add_paragraph()
        p_w1.text = desc
        p_w1.font.size = Pt(9.5)
        p_w1.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 6: Dynamic Adaptation & Feedback Loops
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6)
    add_header(slide6, "Real-Time Adaptive Learning Engine", "CONTINUOUS PERSONALIZATION")

    adaptations = [
        ("📉 Too Difficult Rating", "Engine injects intermediate sub-topics and supplementary beginner-friendly exercises into the roadmap before advancing.", RGBColor(239, 68, 68)),
        ("⚡ Too Easy Rating", "Automatically condenses the module, marks foundational milestones as completed, and accelerates path to advanced projects.", COLOR_SUCCESS),
        ("⏭️ Skip Topic Trigger", "Flags item as bypassed, updates prerequisite dependency graph, and re-optimizes remaining roadmap sequence.", COLOR_WARNING),
        ("⏱️ Pacing / Bandwidth Change", "Learner alters available weekly hours -> Engine re-computes target completion dates and milestone deliverables.", COLOR_PRIMARY),
    ]

    for idx, (atitle, adesc, acolor) in enumerate(adaptations):
        row = idx // 2
        col = idx % 2
        c_left = Inches(0.8 + col * 6.0)
        c_top = Inches(1.6 + row * 2.6)
        
        add_card(slide6, c_left, c_top, Inches(5.7), Inches(2.35), bg_color=COLOR_CARD_BG, border_color=acolor)
        
        tb = slide6.shapes.add_textbox(c_left + Inches(0.25), c_top + Inches(0.2), Inches(5.2), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = atitle
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = acolor
        
        p_d = tf.add_paragraph()
        p_d.text = f"\n{adesc}"
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 7: AI Mentor & Dynamic Assessments
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7)
    add_header(slide7, "Grounded AI Mentor & Skill Assessments", "CONVERSATIONAL & ASSESSMENT AGENTS")

    # Left: AI Mentor Card
    add_card(slide7, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), bg_color=COLOR_CARD_BG, border_color=COLOR_PRIMARY)
    tb_m = slide7.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.0), Inches(4.8))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True
    
    p = tf_m.paragraphs[0]
    p.text = "💬 24/7 Context-Grounded AI Mentor"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    mentor_pts = [
        "Fully Grounded in Student State: Injects current progress (e.g. 28.5%), target role (e.g. Backend Java Dev), and skill gaps into prompt context.",
        "Zero Hallucination of Progress: Never assumes learner status—reads live SQLite DB records.",
        "Actionable Guidance: Explains complex topics, recommends specific next roadmap steps, and debugs code snippets.",
        "Hybrid AI Provider: Uses Google Gemini models with full deterministic offline fallback."
    ]
    for pt in mentor_pts:
        p_pt = tf_m.add_paragraph()
        p_pt.text = f"\n• {pt}"
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = COLOR_TEXT_MUTED

    # Right: Assessments Card
    add_card(slide7, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2), bg_color=COLOR_CARD_BG, border_color=COLOR_SUCCESS)
    tb_a = slide7.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.0), Inches(4.8))
    tf_a = tb_a.text_frame
    tf_a.word_wrap = True
    
    p = tf_a.paragraphs[0]
    p.text = "📝 Dynamic AI Skill Assessments"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUCCESS

    assess_pts = [
        "Targeted Topic Quizzes: Generates multi-choice and scenario questions matched to current roadmap module.",
        "Instant Grading & Explanations: Explains why choices are right/wrong to reinforce deep concept mastery.",
        "Automated Skill Score Updating: Passing an assessment directly updates the learner's skill gap rating (e.g. Core Java 3.5 -> 4.5/5).",
        "Adaptive Difficulty: Questions scale in complexity based on previous test scores."
    ]
    for pt in assess_pts:
        p_pt = tf_a.add_paragraph()
        p_pt.text = f"\n• {pt}"
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 8: Tech Stack & Engineering Highlights
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8)
    add_header(slide8, "Technical Stack & Engineering Highlights", "FULL-STACK ARCHITECTURE")

    tech_cols = [
        ("Frontend Technologies", [
            "React 19 with TypeScript",
            "Vite (Lightning Fast HMR)",
            "Lucide React Icons",
            "Radar & Progress Visualizations",
            "Tailored CSS Theme & Glassmorphism",
            "Modular Component Architecture"
        ], COLOR_PRIMARY),
        ("Backend Technologies", [
            "FastAPI Asynchronous Framework",
            "SQLAlchemy 2.0 Async ORM",
            "AioSQLite Database Engine",
            "Pydantic V2 Schema Validation",
            "Uvicorn ASGI Server",
            "HTTPX Async Client for AI APIs"
        ], COLOR_ACCENT),
        ("AI & Recommendation Core", [
            "Google Gemini AI SDK & REST",
            "Zero-Key Deterministic Fallback",
            "NLP Entity & Skill Extraction",
            "8-Factor Weighted Scoring Engine",
            "Dynamic Quiz Generation Engine",
            "Prerequisite DAG Validation"
        ], COLOR_SUCCESS)
    ]

    for idx, (ttitle, titems, tcolor) in enumerate(tech_cols):
        c_left = Inches(0.8 + idx * 4.0)
        add_card(slide8, c_left, Inches(1.6), Inches(3.7), Inches(5.2), bg_color=COLOR_CARD_BG, border_color=tcolor)
        
        tb = slide8.shapes.add_textbox(c_left + Inches(0.2), Inches(1.8), Inches(3.3), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = ttitle
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = tcolor
        
        for item in titems:
            p_i = tf.add_paragraph()
            p_i.text = f"\n✔ {item}"
            p_i.font.size = Pt(11)
            p_i.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 9: Interactive Demo & Alex Chen Persona
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide9)
    add_header(slide9, "Live Demo Case Study: Alex Chen Profile", "DEMO EXPERIENCE")

    # Left: Persona Card
    add_card(slide9, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), bg_color=COLOR_CARD_BG, border_color=COLOR_HIGHLIGHT)
    tb_p = slide9.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.0), Inches(4.8))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True
    
    p = tf_p.paragraphs[0]
    p.text = "👤 Pre-Seeded Learner: Alex Chen"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_HIGHLIGHT

    persona_info = [
        ("Target Career Role", "Backend Java Developer (Enterprise)"),
        ("Current Skill Baseline", "Core Java (3.5/5), SQL (3.0/5), OOP (3.5/5)"),
        ("Critical Skill Gaps", "Spring Boot (1.0/5), Microservices (0.5/5), Redis"),
        ("Weekly Commitment", "10 Hours / week with target date in 6 months"),
        ("Current Progress", "28.5% Completed (3 of 12 roadmap items finished)"),
        ("Active Streak", "12 Consecutive Days · 34.5 Total Hours Learned")
    ]
    for label, val in persona_info:
        p1 = tf_p.add_paragraph()
        p1.text = f"\n• {label}: "
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_LIGHT
        p1.text += val

    # Right: Dashboard & Roadmap highlights
    add_card(slide9, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2), bg_color=COLOR_CARD_BG, border_color=COLOR_SUCCESS)
    tb_d = slide9.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.0), Inches(4.8))
    tf_d = tb_d.text_frame
    tf_d.word_wrap = True
    
    p = tf_d.paragraphs[0]
    p.text = "📊 Interactive Workflow in Demo"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUCCESS

    demo_steps = [
        "1. Instant Demo Login: One-click load bypassing cold-start onboarding.",
        "2. Radar Visualization: View live gap comparison vs. industry target requirements.",
        "3. Interactive Roadmap: Toggle item statuses, view prerequisites, milestones, and projects.",
        "4. Live Adaptive Feedback: Rate a module 'Too Hard' to watch the curriculum adapt.",
        "5. AI Mentor Chat: Ask questions like 'Why do I need Redis caching?' and receive tailored answers.",
        "6. Skill Quiz Studio: Take a quiz on Spring Boot to boost mastery scores."
    ]
    for step in demo_steps:
        p_s = tf_d.add_paragraph()
        p_s.text = f"\n{step}"
        p_s.font.size = Pt(11)
        p_s.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 10: Future Scope & Roadmap
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide10)
    add_header(slide10, "Future Innovations & Scalability", "FUTURE ROADMAP")

    future_goals = [
        ("Multi-Agent Collaborative Mentorship", "Specialized AI agents for code review, resume optimization, and mock technical interviews.", COLOR_PRIMARY),
        ("Enterprise & LMS Integrations", "Seamless connectors for Canvas, Moodle, Coursera, GitHub, and corporate HR platforms (Workday).", COLOR_ACCENT),
        ("Interactive Code Sandbox", "In-browser live code execution environment for instant project validation and debugging.", COLOR_HIGHLIGHT),
        ("Peer & Cohort Learning Tracks", "Matching learners with similar paths for peer accountability, group projects, and study sessions.", COLOR_SUCCESS),
        ("Mobile Application (iOS/Android)", "Bite-sized mobile micro-learning assessments, streak notifications, and audio mentor mode.", COLOR_WARNING),
        ("Automated Resume & Portfolio Sync", "Auto-generating verifiable GitHub project badges and LinkedIn skill verification credentials.", COLOR_PRIMARY)
    ]

    for idx, (gtitle, gdesc, gcolor) in enumerate(future_goals):
        row = idx // 3
        col = idx % 3
        c_left = Inches(0.8 + col * 4.0)
        c_top = Inches(1.6 + row * 2.6)
        
        add_card(slide10, c_left, c_top, Inches(3.7), Inches(2.35), bg_color=COLOR_CARD_BG, border_color=gcolor)
        
        tb = slide10.shapes.add_textbox(c_left + Inches(0.2), c_top + Inches(0.2), Inches(3.3), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = gtitle
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = gcolor
        
        p_d = tf.add_paragraph()
        p_d.text = f"\n{gdesc}"
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 11: Summary & Project Links
    # ==========================================
    slide11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide11)
    add_header(slide11, "Project Summary & Resources", "THANK YOU / Q&A")

    add_card(slide11, Inches(1.5), Inches(1.6), Inches(10.333), Inches(5.2), bg_color=COLOR_CARD_BG, border_color=COLOR_PRIMARY)
    
    tb = slide11.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(9.333), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "AI Path Finder — Bridging the Gap from Learner to Professional"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_HIGHLIGHT

    summary_pts = [
        ("GitHub Repository", "https://github.com/arunkumarar45/AI-Path-Finder"),
        ("Key Innovation", "Dynamic, real-time adaptation loop connecting feedback to 8-factor recommendation algorithms"),
        ("AI Strategy", "Hybrid Gemini LLM + Offline deterministic fallback ensuring 100% platform availability"),
        ("Architecture", "Modern React 19 Frontend + High Performance FastAPI Backend + SQLite Async ORM"),
        ("Demo Mode", "Full out-of-the-box interactive persona (Alex Chen) pre-loaded with rich data")
    ]
    for title, val in summary_pts:
        p_pt = tf.add_paragraph()
        p_pt.text = f"\n📌 {title}: "
        p_pt.font.size = Pt(13)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_TEXT_LIGHT
        p_pt.text += val

    p_end = tf.add_paragraph()
    p_end.text = "\nThank You! Open for Questions & Feedback."
    p_end.font.size = Pt(14)
    p_end.font.bold = True
    p_end.font.color.rgb = COLOR_SUCCESS

    # Save presentation
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    out_dir = r"d:\AI Path Finder"
    out_file = os.path.join(out_dir, "AI_Path_Finder_Presentation.pptx")
    create_presentation(out_file)
